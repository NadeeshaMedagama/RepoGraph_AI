# ═══════════════════════════════════════════════════════════════════════════════
# RepoGraph AI - Document Processor
# ═══════════════════════════════════════════════════════════════════════════════
# Processor for document files (DOCX, PDF, PPTX, ODT)
# ═══════════════════════════════════════════════════════════════════════════════

from typing import Set, List
from pathlib import Path

from models import Document, FileCategory, FileType, ProcessingStatus
from processors.base_processor import BaseProcessor
from utils import get_logger

logger = get_logger(__name__)


class DocumentProcessor(BaseProcessor):
    """
    Processor for document files.

    Handles DOCX, PDF, PPTX, and ODT files.
    Extracts text content including paragraphs, tables, and slides.
    """

    @property
    def supported_types(self) -> Set[FileType]:
        return {
            FileType.DOCX,
            FileType.DOC,
            FileType.PDF,
            FileType.PPTX,
            FileType.PPT,
            FileType.ODT,
        }

    @property
    def category(self) -> FileCategory:
        return FileCategory.DOCUMENT

    def extract(self, document: Document) -> Document:
        """
        Extract content from a document file.

        Args:
            document: Document to process

        Returns:
            Document with extracted_text populated
        """
        try:
            document.status = ProcessingStatus.EXTRACTING

            file_type = document.metadata.file_type

            if file_type == FileType.DOCX:
                content = self._extract_docx(document)
            elif file_type == FileType.PDF:
                content = self._extract_pdf(document)
            elif file_type in (FileType.PPTX, FileType.PPT):
                content = self._extract_pptx(document)
            elif file_type == FileType.ODT:
                content = self._extract_odt(document)
            else:
                content = f"Document: {document.metadata.name}"

            document.extracted_text = content
            document.status = ProcessingStatus.COMPLETED

            logger.debug(
                "Document processed",
                file=document.metadata.name,
                type=file_type.value,
                text_length=len(content),
            )

            return document

        except Exception as e:
            document.status = ProcessingStatus.FAILED
            document.error_message = str(e)
            logger.error(f"Failed to process document: {e}")
            return document

    def _extract_docx(self, document: Document) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document as DocxDocument
            from docx.table import Table

            doc = DocxDocument(document.metadata.path)
            texts = [f"Document: {document.metadata.name}\n"]

            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Check if it's a heading
                    if para.style and 'Heading' in para.style.name:
                        texts.append(f"\n## {text}\n")
                    else:
                        texts.append(text)

            # Extract tables
            for table in doc.tables:
                table_text = self._extract_docx_table(table)
                if table_text:
                    texts.append(f"\n[Table]\n{table_text}")

            return "\n".join(texts)

        except Exception as e:
            logger.warning(f"Failed to extract DOCX: {e}")
            return f"Document: {document.metadata.name}\n(Content extraction failed)"

    def _extract_docx_table(self, table) -> str:
        """Extract text from a DOCX table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_pdf(self, document: Document) -> str:
        """Extract text from PDF file."""
        try:
            import pdfplumber

            texts = [f"Document: {document.metadata.name}\n"]

            with pdfplumber.open(document.metadata.path) as pdf:
                for i, page in enumerate(pdf.pages):
                    texts.append(f"\n--- Page {i + 1} ---\n")

                    # Extract text
                    page_text = page.extract_text() or ""
                    if page_text.strip():
                        texts.append(page_text)

                    # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_text = self._format_table(table)
                            texts.append(f"\n[Table]\n{table_text}")

            return "\n".join(texts)

        except Exception as e:
            logger.warning(f"Failed to extract PDF with pdfplumber: {e}")
            return self._extract_pdf_fallback(document)

    def _extract_pdf_fallback(self, document: Document) -> str:
        """Fallback PDF extraction using PyPDF2."""
        try:
            from PyPDF2 import PdfReader

            texts = [f"Document: {document.metadata.name}\n"]

            reader = PdfReader(document.metadata.path)
            for i, page in enumerate(reader.pages):
                texts.append(f"\n--- Page {i + 1} ---\n")
                text = page.extract_text() or ""
                if text.strip():
                    texts.append(text)

            return "\n".join(texts)

        except Exception as e:
            logger.warning(f"Failed to extract PDF: {e}")
            return f"Document: {document.metadata.name}\n(PDF extraction failed)"

    def _extract_pptx(self, document: Document) -> str:
        """Extract text from PowerPoint file."""
        try:
            from pptx import Presentation

            prs = Presentation(document.metadata.path)
            texts = [f"Presentation: {document.metadata.name}\n"]
            texts.append(f"Total slides: {len(prs.slides)}\n")

            for i, slide in enumerate(prs.slides):
                texts.append(f"\n=== Slide {i + 1} ===")

                # Extract title if available
                if slide.shapes.title:
                    texts.append(f"Title: {slide.shapes.title.text}")

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        if shape != slide.shapes.title:
                            texts.append(shape.text)

                    # Extract from tables
                    if shape.has_table:
                        table_text = self._extract_pptx_table(shape.table)
                        if table_text:
                            texts.append(f"\n[Table]\n{table_text}")

            return "\n".join(texts)

        except Exception as e:
            logger.warning(f"Failed to extract PPTX: {e}")
            return f"Presentation: {document.metadata.name}\n(Content extraction failed)"

    def _extract_pptx_table(self, table) -> str:
        """Extract text from a PPTX table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_odt(self, document: Document) -> str:
        """Extract text from ODT file."""
        try:
            from odf import text, load
            from odf.text import P

            doc = load(document.metadata.path)
            texts = [f"Document: {document.metadata.name}\n"]

            all_paras = doc.getElementsByType(P)
            for para in all_paras:
                # Get text content recursively
                para_text = ""
                for node in para.childNodes:
                    if hasattr(node, 'data'):
                        para_text += node.data
                    elif hasattr(node, 'childNodes'):
                        for child in node.childNodes:
                            if hasattr(child, 'data'):
                                para_text += child.data

                if para_text.strip():
                    texts.append(para_text)

            return "\n".join(texts)

        except Exception as e:
            logger.warning(f"Failed to extract ODT: {e}")
            return f"Document: {document.metadata.name}\n(ODT extraction failed)"

    def _format_table(self, table: List[List]) -> str:
        """Format a table as text."""
        if not table:
            return ""

        rows = []
        for row in table:
            cells = [str(cell or "").strip() for cell in row]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
